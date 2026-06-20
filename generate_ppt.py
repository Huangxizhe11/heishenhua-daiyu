#!/usr/bin/env python3
"""Generate hackathon project presentation for 黑神话·林黛玉"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette (matches game config) ──
RED      = RGBColor(0xE9, 0x45, 0x60)
TEAL     = RGBColor(0x4E, 0xCD, 0xC4)
GOLD     = RGBColor(0xFF, 0xD9, 0x3D)
DARK_BG  = RGBColor(0x0D, 0x02, 0x21)
PURPLE   = RGBColor(0x1A, 0x0A, 0x2E)
PINK     = RGBColor(0xFF, 0xB6, 0xC1)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xCC, 0xCC, 0xCC)
GRAY     = RGBColor(0x99, 0x99, 0x99)
DARK_CARD = RGBColor(0x15, 0x05, 0x35)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helpers ──
def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_para(tf, text, font_size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT, space_before=Pt(4), font_name="Microsoft YaHei"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    if space_before:
        p.space_before = space_before
    return p

def add_accent_line(slide, left, top, width, color=RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, items, accent=TEAL, icon=""):
    card = add_shape(slide, left, top, width, height, fill_color=DARK_CARD, line_color=accent, line_width=Pt(1.5))
    # title
    t = f"{icon}  {title}" if icon else title
    add_text_box(slide, left + Inches(0.25), top + Inches(0.15), width - Inches(0.5), Inches(0.5),
                 t, font_size=18, color=accent, bold=True)
    # items
    y = top + Inches(0.6)
    for item in items:
        add_text_box(slide, left + Inches(0.25), y, width - Inches(0.5), Inches(0.35),
                     f"•  {item}", font_size=13, color=LIGHT)
        y += Inches(0.32)
    return card

def add_number_badge(slide, left, top, number, color=RED):
    size = Inches(0.6)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)

# ══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# decorative top bar
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=RED)

# Subtitle tag
add_shape(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(0.5), fill_color=PURPLE, line_color=RED, line_width=Pt(1))
add_text_box(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(0.5),
             "腾讯黑客松  参赛作品", font_size=16, color=RED, bold=True, align=PP_ALIGN.CENTER)

# Main title
add_text_box(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.2),
             "黑 神 话 · 林 黛 玉", font_size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# English subtitle
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.6),
             "Black Myth: Lin Daiyu", font_size=24, color=TEAL, bold=False, align=PP_ALIGN.CENTER)

# Tagline
add_accent_line(slide, Inches(5.5), Inches(4.3), Inches(2.3), RED)
add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.6),
             "潇湘妃子，降世除魔", font_size=22, color=GOLD, bold=False, align=PP_ALIGN.CENTER)

# Bottom info
add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.5),
             "《黑神话：悟空》 × 《红楼梦》  ·  3D 动作角色扮演浏览器游戏", font_size=16, color=LIGHT, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.5),
             "Three.js + WebAudio API  ·  全程序化生成  ·  零外部模型", font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

# decorative bottom bar
add_shape(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=TEAL)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: CONCEPT / MOTIVATION
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=RED)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
             "创意来源", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(1.5), RED)

# Left: two source works
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), fill_color=DARK_CARD, line_color=RED, line_width=Pt(1))
add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.5),
             "两大中国文化 IP 的碰撞", font_size=20, color=GOLD, bold=True)
txBox = add_text_box(slide, Inches(1.1), Inches(2.1), Inches(5), Inches(1.8),
                     "", font_size=15, color=LIGHT)
tf = txBox.text_frame
tf.paragraphs[0].text = "🐵  《黑神话：悟空》— 顶级动作游戏的战斗体验"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.color.rgb = LIGHT
tf.paragraphs[0].font.name = "Microsoft YaHei"
add_para(tf, "📕  《红楼梦》— 中国古典文学巅峰之作", 15, LIGHT)
add_para(tf, "", 8, LIGHT)
add_para(tf, "将黛玉葬花化为战斗技能，将金陵十二简化为BOSS", 15, TEAL, True)

# Right: the fusion
add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.5), fill_color=DARK_CARD, line_color=TEAL, line_width=Pt(1))
add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5), Inches(0.5),
             "化学反应", font_size=20, color=TEAL, bold=True)
txBox = add_text_box(slide, Inches(7.1), Inches(2.1), Inches(5), Inches(1.8),
                     "", font_size=15, color=LIGHT)
tf = txBox.text_frame
tf.paragraphs[0].text = "▸  葬花锄 → 葬花剑（近战武器）"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.color.rgb = LIGHT
tf.paragraphs[0].font.name = "Microsoft YaHei"
add_para(tf, "▸  黛玉眼泪 → 泪雨术（远程法术）", 15, LIGHT)
add_para(tf, "▸  颦颦一笑 → 治疗 + 群体控制", 15, LIGHT)
add_para(tf, "▸  绛珠仙草 → 天魁星（终极技能）", 15, LIGHT)
add_para(tf, "▸  金锁 · 妒火 · 风月宝镜 → 三大BOSS机制", 15, GOLD, True)

# Bottom: story premise
add_shape(slide, Inches(0.8), Inches(4.4), Inches(11.7), Inches(2.6), fill_color=DARK_CARD, line_color=GOLD, line_width=Pt(1))
add_text_box(slide, Inches(1.1), Inches(4.5), Inches(11), Inches(0.5),
             "故事背景", font_size=20, color=GOLD, bold=True)
txBox = add_text_box(slide, Inches(1.1), Inches(5.0), Inches(11), Inches(1.8),
                     "", font_size=15, color=LIGHT)
tf = txBox.text_frame
tf.paragraphs[0].text = "仙境崩塌，妖魔横行。绛珠仙草转世的林黛玉，手持花锄化成的神器，踏上降妖除魔之路。"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = LIGHT
tf.paragraphs[0].font.name = "Microsoft YaHei"
add_para(tf, "穿越潇湘馆、荣国府、太虚幻境三大场景，击败薛宝钗、赵姨娘、镜中魔三大BOSS，", 16, LIGHT)
add_para(tf, "拯救崩坏的红楼仙境。", 16, RED, True)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: GAMEPLAY & FEATURES
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=RED)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
             "核心玩法", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(1.5), RED)

# Combat system card
add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "⚔️  战斗系统", [
             "连击系统：连续攻击伤害递增 10%",
             "闪避翻滚：无敌帧 + 冷却时间",
             "MP 技能系统：释放技能消耗MP",
             "完美闪避 / 精确打击判定",
             "低血量心跳视觉反馈",
         ], accent=RED, icon="")

# Skills card
add_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "🌟  四大技能", [
             "左键 · 葬花剑 — 基础近战连击",
             "右键 · 泪雨术 — 远程花雨弹幕",
             "Q · 天魁星 — 终极爆发（400伤害）",
             "E · 颦颦一笑 — 治疗+眩晕控制",
             "Space · 翻滚闪避（无敌帧）",
         ], accent=TEAL, icon="")

# Visual feedback card
add_card(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(2.8),
         "✨  视觉反馈", [
             "伤害数字浮动（普通/暴击/治疗）",
             "屏幕震动 + 低血量红色心跳",
             "连击计数器实时显示",
             "BOSS阶段转换全屏通告",
             "花瓣粒子全程飘落",
         ], accent=GOLD, icon="")

# Boss system - bottom section
add_shape(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.6), fill_color=DARK_CARD, line_color=RED, line_width=Pt(1.5))
add_text_box(slide, Inches(1.1), Inches(4.7), Inches(10), Inches(0.5),
             "👹  三大 BOSS — 每个BOSS独有机制 + 三阶段变身", font_size=20, color=RED, bold=True)

# Boss 1
add_number_badge(slide, Inches(1.2), Inches(5.35), 1, GOLD)
add_text_box(slide, Inches(1.9), Inches(5.3), Inches(3.2), Inches(0.4),
             "薛宝钗 · 金锁迷障", font_size=16, color=GOLD, bold=True)
add_text_box(slide, Inches(1.9), Inches(5.7), Inches(3.2), Inches(1.2),
             "防御反击型\n寒气吐息(减速) + 牡丹绽放(定身)\n金锁格挡反弹 · HP: 5000", font_size=12, color=LIGHT)

# Boss 2
add_number_badge(slide, Inches(5.2), Inches(5.35), 2, RED)
add_text_box(slide, Inches(5.9), Inches(5.3), Inches(3.2), Inches(0.4),
             "赵姨娘 · 妒火炼狱", font_size=16, color=RED, bold=True)
add_text_box(slide, Inches(5.9), Inches(5.7), Inches(3.2), Inches(1.2),
             "远程召唤型\n纸人追踪 + 火焰吐息(多段)\n诅咒弹幕 · HP: 7000", font_size=12, color=LIGHT)

# Boss 3
add_number_badge(slide, Inches(9.2), Inches(5.35), 3, RGBColor(0x9B, 0x59, 0xB6))
add_text_box(slide, Inches(9.9), Inches(5.3), Inches(3.2), Inches(0.4),
             "镜中魔 · 太虚幻境", font_size=16, color=RGBColor(0x9B, 0x59, 0xB6), bold=True)
add_text_box(slide, Inches(9.9), Inches(5.7), Inches(3.2), Inches(1.2),
             "瞬移分身型\n镜片弹幕 + 瞬移斩击\n分身系统 + 碎裂AOE · HP: 10000", font_size=12, color=LIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: TECH ARCHITECTURE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=TEAL)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
             "技术架构", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(1.5), TEAL)

# Tech stack cards - top row
add_card(slide, Inches(0.8), Inches(1.5), Inches(2.7), Inches(2.4),
         "Three.js r128", [
             "WebGL 3D 渲染引擎",
             "程序化几何体建模",
             "无外部模型/贴图",
             "实时光照与粒子系统",
         ], accent=TEAL, icon="🧊")

add_card(slide, Inches(3.8), Inches(1.5), Inches(2.7), Inches(2.4),
         "WebAudio API", [
             "振荡器程序化音效",
             "AI 生成 BGM (MiniMax)",
             "5首原创配乐",
             "实时音效混合",
         ], accent=RED, icon="🔊")

add_card(slide, Inches(6.8), Inches(1.5), Inches(2.7), Inches(2.4),
         "纯前端技术", [
             "零构建工具 / 零框架",
             "单 HTML 文件入口",
             "原生 ES6+ JavaScript",
             "CSS3 动画与布局",
         ], accent=GOLD, icon="⚡")

add_card(slide, Inches(9.8), Inches(1.5), Inches(2.7), Inches(2.4),
         "AI 辅助开发", [
             "MiniMax 生成 BGM 音乐",
             "AI 生成宣传海报",
             "Claude 辅助编码",
             "5小时完成全流程",
         ], accent=PINK, icon="🤖")

# Architecture diagram - bottom
add_shape(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(3.0), fill_color=DARK_CARD, line_color=TEAL, line_width=Pt(1))
add_text_box(slide, Inches(1.1), Inches(4.3), Inches(10), Inches(0.5),
             "📦  模块化架构 — 8 个 JS 模块，各司其职", font_size=18, color=TEAL, bold=True)

modules = [
    ("config.js", "游戏配置", "玩家/BOSS数值、颜色常量"),
    ("levels.js", "关卡定义", "3关卡配置、剧情文本、掉落表"),
    ("audio.js", "音频管理", "WebAudio音效 + MP3 BGM"),
    ("world.js", "世界生成", "地形、天空、光照、粒子"),
    ("player.js", "玩家系统", "黛玉模型、移动、战斗、状态"),
    ("boss.js", "BOSS系统", "3个BOSS模型 + 独立AI"),
    ("vfx.js", "视觉特效", "粒子、弹幕、屏幕效果"),
    ("main.js", "游戏主循环", "状态机、相机、UI绑定"),
]

x_start = Inches(1.0)
y_start = Inches(4.9)
card_w = Inches(2.7)
card_h = Inches(1.0)
gap_x = Inches(0.15)
gap_y = Inches(0.15)

for i, (name, label, desc) in enumerate(modules):
    col = i % 4
    row = i // 4
    x = x_start + col * (card_w + gap_x)
    y = y_start + row * (card_h + gap_y)
    c = [TEAL, RED, GOLD, PINK, TEAL, RED, GOLD, PINK][i]
    add_shape(slide, x, y, card_w, card_h, fill_color=PURPLE, line_color=c, line_width=Pt(1))
    add_text_box(slide, x + Inches(0.1), y + Inches(0.05), card_w - Inches(0.2), Inches(0.35),
                 name, font_size=13, color=c, bold=True)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.35), card_w - Inches(0.2), Inches(0.3),
                 label, font_size=12, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.6), card_w - Inches(0.2), Inches(0.35),
                 desc, font_size=10, color=GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: CHARACTER DESIGN
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=PINK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
             "角色设计", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(1.5), PINK)

# Player card
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.5), fill_color=DARK_CARD, line_color=PINK, line_width=Pt(1.5))
add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5.3), Inches(0.5),
             "🌸  林黛玉 — 主角", font_size=22, color=PINK, bold=True)

player_features = [
    ("外观设计", "紫色汉服（4层飘逸裙摆）· 金色腰带 · 发髻+5缕飘发 · 双飘带"),
    ("武器", "花锄神器 — 竹柄 + 红光刀刃 + 花朵装饰"),
    ("表情系统", "忧虑眉 · 战斗表情变化"),
    ("氛围效果", "发光光环环 · 花瓣粒子跟随"),
    ("动作系统", "行走 · 奔跑 · 呼吸动画 · 攻击动作 · 翻滚"),
    ("状态系统", "减速（宝钗寒气）· 定身（牡丹绽放）"),
]
y = Inches(2.2)
for label, desc in player_features:
    add_text_box(slide, Inches(1.3), y, Inches(1.5), Inches(0.35),
                 f"▸ {label}", font_size=13, color=TEAL, bold=True)
    add_text_box(slide, Inches(2.8), y, Inches(3.5), Inches(0.55),
                 desc, font_size=12, color=LIGHT)
    y += Inches(0.48)

# Boss cards - right side
bosses = [
    ("👑  薛宝钗", GOLD, "金色铠甲 · 流苏裙摆\n金锁/盾牌光环 · 牡丹纹饰\n防御反击型 · 冰系技能"),
    ("🔥  赵姨娘", RED, "火焰主题 · 橙红配色\n火焰特效 · 纸人召唤\n远程召唤型 · 火系技能"),
    ("🪞  镜中魔", RGBColor(0x9B, 0x59, 0xB6), "紫罗兰色 · 镜面碎片效果\n分身能力 · 瞬移特效\n瞬移分身型 · 幻系技能"),
]

for i, (name, color, desc) in enumerate(bosses):
    y = Inches(1.5) + i * Inches(1.9)
    add_shape(slide, Inches(7.0), y, Inches(5.5), Inches(1.7), fill_color=DARK_CARD, line_color=color, line_width=Pt(1.5))
    add_text_box(slide, Inches(7.3), y + Inches(0.1), Inches(5), Inches(0.4),
                 name, font_size=18, color=color, bold=True)
    add_text_box(slide, Inches(7.3), y + Inches(0.5), Inches(5), Inches(1.0),
                 desc, font_size=12, color=LIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: VISUAL EFFECTS & WORLD
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
             "视觉与世界", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(1.5), GOLD)

# Color palette
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.5), fill_color=DARK_CARD, line_color=GOLD, line_width=Pt(1))
add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.5),
             "🎨  色彩体系 — 暗黑国风", font_size=18, color=GOLD, bold=True)

colors = [
    ("#E94560", "攻击红", RED),
    ("#4ECDC4", "治愈青", TEAL),
    ("#FFD93D", "宝藏金", GOLD),
    ("#1A0A2E", "神秘紫", RGBColor(0x1A, 0x0A, 0x2E)),
    ("#FFB6C1", "花瓣粉", PINK),
    ("#87CEEB", "泪珠蓝", RGBColor(0x87, 0xCE, 0xEB)),
]
x = Inches(1.2)
for hex_code, label, rgb in colors:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(0.7), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()
    add_text_box(slide, x - Inches(0.05), Inches(3.0), Inches(0.8), Inches(0.3),
                 label, font_size=10, color=LIGHT, align=PP_ALIGN.CENTER)
    x += Inches(0.9)

# VFX features
add_card(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.5),
         "✨  视觉特效系统", [
             "花瓣粒子系统 — 全程飘落的花瓣",
             "剑气波 — 葬花剑攻击弹道",
             "火焰拖尾 — 赵姨娘火焰特效",
             "镜面碎片弹幕 — 镜中魔攻击特效",
             "屏幕震动 + 低血量红色心跳环",
             "伤害数字浮动（普通/暴击/治疗/Boss）",
         ], accent=GOLD, icon="")

# Three levels world
add_shape(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.9), fill_color=DARK_CARD, line_color=TEAL, line_width=Pt(1))
add_text_box(slide, Inches(1.1), Inches(4.4), Inches(10), Inches(0.5),
             "🏔️  三大场景 — 程序化 3D 世界", font_size=18, color=TEAL, bold=True)

levels_info = [
    ("第一关 · 潇湘馆", TEAL, "竹林环绕的潇湘仙境\n紫色与青色主调\n花瓣飘落 · 薄雾弥漫", "薛宝钗"),
    ("第二关 · 荣国府", RED, "宏伟的古典宫殿建筑\n红色与金色主调\n火焰特效 · 庄严肃穆", "赵姨娘"),
    ("第三关 · 太虚幻境", RGBColor(0x9B, 0x59, 0xB6), "虚幻的镜像空间\n紫色与银色主调\n镜面反射 · 时空扭曲", "镜中魔"),
]

for i, (name, color, desc, boss) in enumerate(levels_info):
    x = Inches(1.0) + i * Inches(3.9)
    add_shape(slide, x, Inches(5.0), Inches(3.6), Inches(2.0), fill_color=PURPLE, line_color=color, line_width=Pt(1))
    add_text_box(slide, x + Inches(0.15), Inches(5.1), Inches(3.3), Inches(0.4),
                 name, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), Inches(5.45), Inches(3.3), Inches(1.0),
                 desc, font_size=11, color=LIGHT)
    add_text_box(slide, x + Inches(0.15), Inches(6.45), Inches(3.3), Inches(0.35),
                 f"BOSS: {boss}", font_size=12, color=GOLD, bold=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: DEVELOPMENT TIMELINE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=TEAL)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
             "开发历程", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(1.5), TEAL)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "⏱️  5 小时 · 14 次提交 · 从零到完整游戏", font_size=18, color=GOLD, bold=True)

timeline = [
    ("17:02", "v2 · v3", "3关卡完整Demo + 战斗手感调优 + BOSS AI", TEAL),
    ("17:12", "v4", "角色模型全面重做 + 行走/呼吸动画", RED),
    ("17:19", "v5", "模型迭代 + BGM修复 + 移动方向修正", GOLD),
    ("17:50", "v6", "MiniMax AI 真人BGM + UI布局优化（黑神话风格）", PINK),
    ("18:08", "fix", "BGM连续播放 + 剑光特效 + 代码清理", TEAL),
    ("20:08", "feat", "三BOSS独特外观设计 + 关卡流程修复", RED),
    ("21:38", "feat", "第三关BGM重做 — 唢呐二胡太鼓黑神话红楼梦风格", GOLD),
    ("22:26", "final", "三BOSS阶段性战斗系统 + 按键残留修复 + 音频优化", RGBColor(0x9B, 0x59, 0xB6)),
]

# Timeline line
add_shape(slide, Inches(2.2), Inches(2.1), Pt(3), Inches(5.0), fill_color=TEAL)

for i, (time, ver, desc, color) in enumerate(timeline):
    y = Inches(2.1) + i * Inches(0.6)
    # dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.1), y + Inches(0.05), Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    # time
    add_text_box(slide, Inches(0.8), y, Inches(1.2), Inches(0.3),
                 time, font_size=13, color=color, bold=True, align=PP_ALIGN.RIGHT)
    # version
    add_text_box(slide, Inches(2.5), y, Inches(0.8), Inches(0.3),
                 ver, font_size=12, color=color, bold=True)
    # description
    add_text_box(slide, Inches(3.4), y, Inches(9), Inches(0.3),
                 desc, font_size=13, color=LIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 8: HIGHLIGHTS & SUMMARY
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
             "项目亮点", font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(1.5), GOLD)

highlights = [
    ("🎭", "文化融合", "黑神话 × 红楼梦\n两大国民IP创意碰撞", GOLD),
    ("⚡", "极致效率", "5小时从零到完整\n14次Git提交迭代", TEAL),
    ("🧊", "纯前端3D", "Three.js程序化建模\n零外部模型/贴图", RED),
    ("🤖", "AI 协作", "MiniMax生成BGM\nClaude辅助编码", PINK),
]

for i, (icon, title, desc, color) in enumerate(highlights):
    x = Inches(0.8) + i * Inches(3.1)
    add_shape(slide, x, Inches(1.5), Inches(2.8), Inches(2.2), fill_color=DARK_CARD, line_color=color, line_width=Pt(1.5))
    add_text_box(slide, x, Inches(1.6), Inches(2.8), Inches(0.6),
                 icon, font_size=36, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.2), Inches(2.8), Inches(0.4),
                 title, font_size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.4), Inches(0.8),
                 desc, font_size=13, color=LIGHT, align=PP_ALIGN.CENTER)

# Key stats
add_shape(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(1.5), fill_color=DARK_CARD, line_color=GOLD, line_width=Pt(1))
stats = [
    ("8", "JS 模块"),
    ("3", "关卡场景"),
    ("3", "BOSS 战"),
    ("4", "玩家技能"),
    ("5", "首BGM"),
    ("~216KB", "代码量"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(1.2) + i * Inches(1.9)
    add_text_box(slide, x, Inches(4.1), Inches(1.5), Inches(0.7),
                 num, font_size=36, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(4.7), Inches(1.5), Inches(0.4),
                 label, font_size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# Quote
add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.4), fill_color=PURPLE, line_color=PINK, line_width=Pt(1))
add_text_box(slide, Inches(1.5), Inches(5.9), Inches(10.3), Inches(0.5),
             "「 花谢花飞花满天，红消香断有谁怜 」", font_size=22, color=PINK, bold=False, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.5),
             "— 黛玉葬花  ·  以诗词入战斗，以文化铸游戏", font_size=14, color=LIGHT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 9: THANK YOU
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), fill_color=RED)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.2),
             "感 谢 聆 听", font_size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(3.3), Inches(2.3), RED)

add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.6),
             "黑神话·林黛玉  —  潇湘妃子，降世除魔", font_size=22, color=GOLD, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.5),
             "Three.js + WebAudio API  ·  纯前端 3D 动作游戏", font_size=16, color=TEAL, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.5),
             "腾讯黑客松 2026  ·  所有图形与音效均为程序化生成", font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Bottom decorative bar
add_shape(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=TEAL)

# ── Save ──
output = "/Users/admin/Documents/project/腾讯黑客松项目/黑神话daiyu/黑神话林黛玉-项目介绍.pptx"
prs.save(output)
print(f"✅ PPT saved: {output}")
